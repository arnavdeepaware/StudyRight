"""
Document Extraction and Processing Module

This module provides functionality to:
1. Extract text from various document formats (PDF, DOCX, TXT, PPT, PPTX)
   and images (PNG, JPEG)
2. Process extracted text with Google's Gemini AI
3. Save the analysis results as JSON in the prompts directory

The processed output includes core concepts, explanations, visual descriptions,
and voiceover scripts suitable for educational video creation.

Dependencies:
- vertexai: For communicating with Google's Gemini AI
- PyPDF2: For PDF text extraction
- python-docx: For Word document text extraction
- python-pptx: For PowerPoint text extraction
- pytesseract: For image text extraction (OCR)
- Pillow: For image processing

Usage:
- To process a specific file: process_uploaded_file("filename.pdf")
- To process the oldest file in uploads: process_oldest_file()
- To extract file info without processing: extract_file_info("filename.pdf")
"""

from vertexai.generative_models import GenerativeModel
import vertexai
import os
import json
import PyPDF2
import docx
import datetime
import pytesseract
from PIL import Image
from google.oauth2 import service_account
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

def extract_text_from_file(filepath):
    """
    Extract text from different file types
    
    Args:
        filepath (str): Path to the file
        
    Returns:
        str: Extracted text or error message
    """
    # Text files
    if filepath.endswith('.txt'):
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            return f"Error extracting text file: {str(e)}"
    
    # PDF files
    elif filepath.endswith('.pdf'):
        text = ""
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    text += pdf_reader.pages[page_num].extract_text() + "\n"
            return text if text.strip() else "PDF contained no extractable text (might be scanned)"
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"
    
    # Word documents
    elif filepath.endswith('.docx'):
        try:
            doc = docx.Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text if text.strip() else "Word document contained no extractable text"
        except Exception as e:
            return f"Error extracting DOCX text: {str(e)}"
    
    # PowerPoint presentations
    elif filepath.endswith('.ppt') or filepath.endswith('.pptx'):
        try:
            from pptx import Presentation
            
            # For .pptx files
            if filepath.endswith('.pptx'):
                presentation = Presentation(filepath)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text if text.strip() else "PowerPoint contained no extractable text"
            
            # For .ppt files (legacy format)
            else:
                # PPT files require conversion or specialized libraries
                return "Legacy PowerPoint (.ppt) files have limited text extraction. Consider converting to .pptx"
                
        except ImportError:
            return "PowerPoint extraction requires python-pptx package. Install with: pip install python-pptx"
        except Exception as e:
            return f"Error extracting PowerPoint text: {str(e)}"
    
    # Image files (PNG, JPEG)
    elif filepath.lower().endswith(('.png', '.jpg', '.jpeg')):
        try:
            # Check if pytesseract is properly installed
            image = Image.open(filepath)
            text = pytesseract.image_to_string(image)
            return text if text.strip() else "No text detected in image"
        except ImportError:
            return "Image text extraction requires pytesseract. Install with: pip install pytesseract"
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    # Unsupported file type
    return f"Unsupported file type: {os.path.splitext(filepath)[1]}"

def extract_file_info(filename):
    """
    Extract information from a file without processing it with Gemini
    
    Args:
        filename (str): Name of the file in the uploads directory
        
    Returns:
        dict: File information including extracted text and metadata
    """
    # Construct file path
    uploads_dir = 'uploads'
    filepath = os.path.join(uploads_dir, filename)
    
    # Check if file exists
    if not os.path.exists(filepath):
        return {"error": f"File {filename} not found in uploads directory"}
    
    # Get file stats
    file_stats = os.stat(filepath)
    file_size = file_stats.st_size
    mod_time = datetime.datetime.fromtimestamp(file_stats.st_mtime).isoformat()
    
    # Extract text from file
    extracted_text = extract_text_from_file(filepath)
    
    # Create and return the file info dictionary
    file_info = {
        "filename": filename,
        "filepath": filepath,
        "file_size_bytes": file_size,
        "last_modified": mod_time,
        "extraction_time": datetime.datetime.now().isoformat(),
        "extracted_text": extracted_text,
        "text_preview": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
    }
    
    return file_info

def send_to_gemini(text):
    """
    Send extracted text to Google's Gemini AI for processing
    
    Args:
        text (str): Text content to be analyzed
        
    Returns:
        str: Gemini's analysis or error message
    """
    try:
        # Use service account credentials
        credentials = service_account.Credentials.from_service_account_file(
            'vertexai_key.json'
        )

        # Get project details from .env file
        PROJECT_ID = os.getenv('PROJECT_ID').strip('"').strip()
        LOCATION = os.getenv('LOCATION').strip('"').strip()

        # Initialize Vertex AI
        vertexai.init(
            project=PROJECT_ID,
            location=LOCATION,
            credentials=credentials
        )

        # Use Gemini 2.0 Flash model
        model = GenerativeModel("gemini-2.0-flash")

        
        prompt_template = '''
You are assisting a college student in turning their academic material into short, fun, and engaging social media-style videos.

Format exactly 5-6 key concepts like this (no extra newlines or spaces):

**1. Title: [Title]**
**Caption:** [Caption with hashtags]
**Visual:** [Visual description]
**Voiceover:** "[Conversational script]"

**2. Title: [Title]**
**Caption:** [Caption]
**Visual:** [Visual]
**Voiceover:** "[Script]"

Content to process:
'''
        
        # Format the prompt with the extracted text
        prompt = f"{prompt_template}\n\n{text[:5000]}"

        response = model.generate_content(prompt)
        return response.text

    except Exception as e:
        return f"Error with Gemini API: {str(e)}"

def process_extracted_info(file_info):
    """
    Process already extracted file information with Gemini and save results
    
    Args:
        file_info (dict): File information dictionary from extract_file_info
        
    Returns:
        dict: Processing result summary
    """
    # Check if there was an error during extraction
    if "error" in file_info:
        return file_info
    
    # Process with Gemini
    gemini_response = send_to_gemini(file_info["extracted_text"])
    
    # Create a summary object
    summary = {
        "filename": file_info["filename"],
        "processed_time": datetime.datetime.now().isoformat(),
        "extraction_time": file_info["extraction_time"],
        "extracted_text_preview": file_info["text_preview"],
        "gemini_analysis": gemini_response
    }
    
    # Ensure prompts directory exists
    prompts_dir = 'prompts'
    if not os.path.exists(prompts_dir):
        os.makedirs(prompts_dir)
    
    # Save as JSON in the prompts directory
    json_path = os.path.join(prompts_dir, 'all_prompts.json')
    
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(summary, json_file, indent=2, ensure_ascii=False)
    
    return {
        "message": "File successfully processed",
        "summary_file": json_path,
        "original_file": file_info["filepath"]
    }

def get_files_in_uploads():
    """
    Get list of files in the uploads directory
    
    Returns:
        list: List of valid files for processing (excluding JSON files)
    """
    uploads_dir = 'uploads'
    if not os.path.exists(uploads_dir):
        os.makedirs(uploads_dir)
    
    # Filter out JSON files from processing
    files = [f for f in os.listdir(uploads_dir) if os.path.isfile(os.path.join(uploads_dir, f)) and not f.endswith('.json')]
    
    return files

def get_oldest_file():
    """
    Get the oldest file in the uploads directory based on modification time
    
    Returns:
        str or None: Filename of the oldest file, or None if no files found
    """
    files = get_files_in_uploads()
    
    if not files:
        return None
    
    uploads_dir = 'uploads'
    oldest_file = min(files, key=lambda f: os.path.getmtime(os.path.join(uploads_dir, f)))
    
    return oldest_file

def process_uploaded_file(filename):
    """
    Process a file from the uploads directory and save analysis results as JSON
    
    Args:
        filename (str): Name of the file to process
        
    Returns:
        dict: Processing result summary
    """
    # Extract file information
    file_info = extract_file_info(filename)
    
    # Process the extracted information
    return process_extracted_info(file_info)

def process_oldest_file():
    """
    Process the oldest file in the uploads directory
    
    Returns:
        dict: Processing result or error message
    """
    oldest_file = get_oldest_file()
    
    if not oldest_file:
        return {"error": "No files found in uploads directory"}
    
    print(f"Processing oldest file: {oldest_file}")
    return process_uploaded_file(oldest_file)

if __name__ == "__main__":
    """
    Main execution block: 
    Processes the oldest file in the uploads directory when script is run directly
    """
    result = process_oldest_file()
    print(json.dumps(result, indent=2))